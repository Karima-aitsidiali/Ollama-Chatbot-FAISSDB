# Dans FileProcessor
import os
import hashlib
import json
import PyPDF2
from docx import Document
import io # Nécessaire pour lire depuis des bytes
from pptx import Presentation

from Utilitaire.EDA_Cleaner import TextPipeline, TextCleaner # Assurez-vous que c'est le bon import

class FileProcessor:
    def __init__(self, chunk_size=384, chunk_overlap=96):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.processed_hashes = set()

    def calculate_hash(self, content_string): # Prend une string
        return hashlib.sha256(content_string.encode('utf-8')).hexdigest()

    def read_content_from_bytes(self, base_filename, file_content_bytes):
        """
        Convertit le contenu d'un fichier (en bytes) en string selon son type (déduit de base_filename).
        """
        file_extension = os.path.splitext(base_filename)[1].lower()
        content_str = ""

        try:
            if file_extension == '.txt':
                content_str = file_content_bytes.decode('utf-8') # Suppose UTF-8
            elif file_extension == '.pdf':
                pdf_stream = io.BytesIO(file_content_bytes)
                pdf_reader = PyPDF2.PdfReader(pdf_stream)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text: # S'assurer que du texte a été extrait
                        content_str += page_text + "\n"
            elif file_extension == '.docx':
                docx_stream = io.BytesIO(file_content_bytes)
                doc = Document(docx_stream)
                for para in doc.paragraphs:
                    content_str += para.text + "\n"
            elif file_extension == '.json':
                json_string = file_content_bytes.decode('utf-8')
                data = json.loads(json_string)
                content_str = json.dumps(data, ensure_ascii=False) # Pour le traitement ultérieur
            elif file_extension == '.pptx':
                pptx_stream = io.BytesIO(file_content_bytes)
                prs = Presentation(pptx_stream)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content_str += shape.text + "\n"
            else:
                raise ValueError(f"Type de fichier non pris en charge : {file_extension} pour {base_filename}")

            return content_str.strip()
        except Exception as e:
            # traceback.print_exc() # Pour débogage
            raise Exception(f"Erreur lors de la lecture du contenu des bytes pour {base_filename} (type: {file_extension}) : {str(e)}")

    def split_into_chunks(self, text):
        # ... (votre logique existante, inchangée) ...
        if not text: return []
        chunks = []
        start = 0
        text_length = len(text)
        while start < text_length:
            end = start + self.chunk_size
            chunk = text[start:end].strip()
            if chunk: chunks.append(chunk)
            start += self.chunk_size - self.chunk_overlap
        return chunks

    def process_file(self, base_filename, file_content_bytes): # Modifié ici
        """
        Traite le contenu d'un fichier (en bytes) : lit, calcule hash, vérifie, nettoie, et retourne chunks.
        """
        try:
            # Étape 1: Lire le contenu des bytes en une chaîne de caractères
            content_string = self.read_content_from_bytes(base_filename, file_content_bytes)

            if not content_string:
                print(f"Aucun contenu textuel extrait de {base_filename}.")
                return None, None # Pas de contenu à traiter

            # Étape 2: Calculer le hash du contenu textuel
            file_hash = self.calculate_hash(content_string)

            if file_hash in self.processed_hashes:
                print(f"Le fichier {base_filename} (hash: {file_hash}) a déjà été traité.")
                return None, file_hash # Retourner None pour les chunks, mais le hash pour info

            # Étape 3: Nettoyage du texte (votre pipeline EDA)
            # Assurez-vous que TextCleaner et TextPipeline sont correctement importés et fonctionnent
            # from Utilitaire.EDA_Cleaner import TextPipeline, TextCleaner # Exemple d'import
            cleaner = TextCleaner()
            pipeline = TextPipeline(cleaner)
            cleaned_content = pipeline.process(content_string) # pipeline.process prend une string

            if not cleaned_content:
                print(f"Contenu vide après nettoyage pour {base_filename}.")
                return None, file_hash # Retourner le hash pour marquer comme traité, mais pas de chunks

            # Étape 4: Diviser en chunks
            chunks = self.split_into_chunks(cleaned_content)

            if not chunks:
                print(f"Aucun chunk généré pour {base_filename} après nettoyage et division.")
                # Décider si on ajoute le hash aux processed_hashes même si pas de chunks
                # Pour l'instant, on ne l'ajoute que si des chunks sont générés et traités
                return None, file_hash

            # Si tout s'est bien passé et que des chunks ont été générés :
            self.processed_hashes.add(file_hash)
            return chunks, file_hash

        except Exception as e:
            print(f"Erreur dans FileProcessor.process_file pour {base_filename}: {str(e)}")
            import traceback
            traceback.print_exc()
            # Il est important de propager l'erreur pour qu'elle soit gérée plus haut
            raise # ou return None, None si vous voulez gérer l'erreur ici et ne pas bloquer